# -*- coding: utf-8 -*-
"""Derecepanel — Koç sidebar özellikleri, yatırımcı slayt stiline yakın kart düzeni (python-pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Renk paleti (Yatırımcı HTML + ekran görüntüsü mor tonları)
C_BG = RGBColor(245, 242, 250)
C_CATEGORY = RGBColor(155, 126, 200)
C_TITLE = RGBColor(45, 31, 61)
C_HEADER = RGBColor(91, 58, 140)
C_HEADER_ALT = RGBColor(107, 77, 158)
C_BODY_TEXT = RGBColor(51, 44, 70)
C_MUTED = RGBColor(92, 83, 104)
C_DECOR = RGBColor(232, 224, 245)
C_ORANGE = RGBColor(242, 100, 25)
C_LINE = RGBColor(210, 200, 228)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG


def add_decor_circles(slide):
    """Köşe dekor — yarı saydam hissi için açık mor daireler."""
    circles = [
        (Inches(10.2), Inches(-0.8), Inches(4.2), Inches(4.2)),
        (Inches(-1.2), Inches(5.0), Inches(3.8), Inches(3.8)),
        (Inches(11.5), Inches(4.8), Inches(3.0), Inches(3.0)),
    ]
    for left, top, w, h in circles:
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(left), int(top), int(w), int(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_DECOR
        sh.line.fill.background()


def add_category_title(slide, category: str, title: str):
    cat = slide.shapes.add_textbox(Inches(0.52), Inches(0.32), Inches(12.2), Inches(0.35))
    tf = cat.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = category
    p.font.name = "Calibri"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_CATEGORY

    tit = slide.shapes.add_textbox(Inches(0.52), Inches(0.58), Inches(12.2), Inches(0.85))
    tf2 = tit.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.name = "Calibri"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = C_TITLE
    p2.line_spacing = 1.05


def _set_header_text(shape, num: str, title: str, icon: str = ""):
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    prefix = (icon + " ") if icon else ""
    p.text = f"{prefix}{num}  {title}"
    p.font.name = "Calibri"
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT


def _set_body_text(shape, text: str, small: bool = False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p0 = tf.paragraphs[0]
    p0.text = text
    p0.font.name = "Calibri"
    p0.font.size = Pt(9.5 if small else 10)
    p0.font.color.rgb = C_BODY_TEXT
    p0.line_spacing = 1.15


def add_card(slide, left, top, width, height, num: str, title: str, body: str, icon: str = "", header_rgb=None):
    hdr_rgb = header_rgb or C_HEADER
    hdr_h = Inches(0.44)
    left, top, width, height = int(left), int(top), int(width), int(height)

    # Dış çerçeve (beyaz gövde)
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    outer.adjustments[0] = 0.04
    outer.fill.solid()
    outer.fill.fore_color.rgb = RGBColor(255, 255, 255)
    outer.line.color.rgb = C_LINE
    outer.line.width = Pt(1)

    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, int(hdr_h))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = hdr_rgb
    hdr.line.fill.background()

    # Başlık metni (şekil üstüne textbox — z-order en üstte kalsın)
    tb_hdr = slide.shapes.add_textbox(left, top, width, int(hdr_h))
    _set_header_text(tb_hdr, num, title, icon)

    body_top = top + int(hdr_h)
    body_h = height - int(hdr_h)
    tb_body = slide.shapes.add_textbox(left, body_top, width, body_h)
    _set_body_text(tb_body, body, small=(height < Inches(1.85)))


def add_accent_bar(slide):
    """Üst turuncu ince şerit — marka vurgusu."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, int(SLIDE_W), int(Inches(0.06)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ORANGE
    bar.line.fill.background()


def build_presentation(path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank

    # --- Kapak ---
    s0 = prs.slides.add_slide(blank)
    add_slide_bg(s0)
    add_decor_circles(s0)
    add_accent_bar(s0)
    t1 = s0.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(1.2))
    p = t1.text_frame.paragraphs[0]
    p.text = "Derecepanel"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_TITLE
    p.font.name = "Calibri"
    p2 = t1.text_frame.add_paragraph()
    p2.text = "Koç paneli — sidebar modülleri"
    p2.font.size = Pt(22)
    p2.font.color.rgb = C_MUTED
    p2.space_before = Pt(8)
    p3 = t1.text_frame.add_paragraph()
    p3.text = "Yatırımcı sunumu · 2026"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_CATEGORY
    p3.space_before = Pt(16)

    # Ortak yerleşim: sol sütun genişliği
    m = Inches(0.52)
    big_w = Inches(3.38)
    gap = Inches(0.32)
    col0 = m
    col1 = col0 + big_w + gap
    card_w = (SLIDE_W - col1 - m - Inches(0.25)) / 2
    row_h = Inches(1.78)
    row_gap = Inches(0.14)
    top0 = Inches(1.38)

    def small_cell(col, row):
        x = col1 + col * (card_w + Inches(0.25))
        y = top0 + row * (row_h + row_gap)
        return x, y, card_w, row_h

    big_h = Inches(5.62)

    # --- Slayt 1: Çekirdek ---
    s1 = prs.slides.add_slide(blank)
    add_slide_bg(s1)
    add_decor_circles(s1)
    add_accent_bar(s1)
    add_category_title(s1, "05A  KOÇ PANELİ — ÇEKİRDEK", "Kurum ve Koç İçin Operasyonel Komuta Katmanı")

    add_card(
        s1,
        col0,
        top0,
        big_w,
        big_h,
        "01",
        "Dashboard & Öğrencilerim",
        "Dashboard: YKS geri sayımı, özet kartlar ve günlük giriş noktası; koçun tüm modüllere dağılan işi tek bakışta toparlaması.\n\n"
        "Öğrencilerim: Atanmış öğrenci portföyü, profil ve takip; panel genelinde filtre ve arama için merkezi öğrenci verisi.",
        "📊",
    )

    cards1 = [
        ("02", "Görüşme Odası", "Yeni Görüşme ile not ve özet kaydı; Görüşme Arşivi ile kurumsal hafıza ve denetlenebilir koçluk döngüsü.", "💬"),
        ("03", "Koçlar", "Kurum içi kadro: ekleme, düzenleme, arama ve detay; çok koçlu yapıda kapasite ve sorumluluk görünürlüğü.", "👥"),
        ("04", "Randevular", "Öğrenci–koç ve kurum içi zaman planı; operasyonel koordinasyonu panel içinde tutar.", "📅"),
        ("05", "Pazarlama Asistanı (BETA)", "Kampanya ve iletişim metni üretimine yakın asistan; erken aşamada CAC ve mesaj tutarlılığı deneyi.", "📣"),
        ("06", "YKS Simülasyon (özet)", "Tercih Sihirbazı, YKS Konuları, Net Sihirbazı, Puan Hesaplama: hedef ve sınav matematiğini tek çatıda.", "🎓"),
        ("07", "Denemeler (özet)", "Sonuçtan optiğe kadar ölçüm zinciri; ayrıntılar 05B slaytında.", "📝"),
    ]
    for i, (num, tit, bod, ic) in enumerate(cards1):
        r, c = divmod(i, 2)
        x, y, w, h = small_cell(c, r)
        add_card(s1, x, y, w, h, num, tit, bod, ic)

    # --- Slayt 2: YKS + Deneme ---
    s2 = prs.slides.add_slide(blank)
    add_slide_bg(s2)
    add_decor_circles(s2)
    add_accent_bar(s2)
    add_category_title(s2, "05B  YKS & DENEME OPERASYONU", "Ölçüm ve Sınav Ekonomisini Ürün İçinde Kapatma")

    add_card(
        s2,
        col0,
        top0,
        big_w,
        big_h,
        "08",
        "YKS Simülasyon (alt menü)",
        "• Tercih Sihirbazı: puan, bölge ve hedeflere göre tercih ve senaryo desteği.\n"
        "• YKS Konuları: müfredat ve konu haritası; planlama ile hizalı referans.\n"
        "• Net Sihirbazı: hedef net bandı ve motivasyon; yol haritası hissi.\n"
        "• Puan Hesaplama: güncel ÖSYM parametreleriyle uyumlu hesap mantığı (takvim güncellenince ürün güncellenir).",
        "🎯",
        C_HEADER_ALT,
    )

    cards2 = [
        ("09", "Sonuç Merkezi", "Deneme sonuçlarının toplandığı ve yönetildiği merkez ekran; operasyonun kalbi.", "📋"),
        ("10", "Analiz ve Raporlama (BETA)", "Derinlemesine analiz ve rapor; veri ürünü olarak ölçek ve farklılaşma katmanı.", "📈"),
        ("11", "Global Deneme Takvimi", "Yaygın deneme takvimine hizalama; kurum dışı senkron.", "🌍"),
        ("12", "Kurum Deneme Takvimi", "Kendi iç sınav ve deneme programının planlanması.", "🏢"),
        ("13", "Deneme Sonuçları Yükleme", "Toplu veya dosya ile sonuç girişi; manuel dağınıklığı azaltır.", "⬆"),
        ("14", "Optik & Okutulan", "Optik okuyucu hattı ve okutulan deneme kayıtları; optikten analize köprü.", "🔍"),
    ]
    for i, (num, tit, bod, ic) in enumerate(cards2):
        r, c = divmod(i, 2)
        x, y, w, h = small_cell(c, r)
        add_card(s2, x, y, w, h, num, tit, bod, ic)

    # --- Slayt 3: Program, MR, Kitap, Test ---
    s3 = prs.slides.add_slide(blank)
    add_slide_bg(s3)
    add_decor_circles(s3)
    add_accent_bar(s3)
    add_category_title(s3, "05C  PLAN, DUYGU, KAYNAK", "Haftalık Ritim, MR ve İçerik Dağıtımı")

    add_card(
        s3,
        col0,
        top0,
        big_w,
        big_h,
        "15",
        "Haftalık Program",
        "Haftalık Program Oluşturucu ile etüt/konu/süre tasarımı; Kayıtlı Haftalık Programlar ile şablon ve arşiv; "
        "Öğrenciye program gönder ile planın öğrenci paneline düşmesi ve bildirimle kapanan döngü.",
        "📆",
    )

    cards3 = [
        ("16", "Konu MR", "Konu bazlı ruh hâli ve yorgunluk sinyalleri; öğrenme verimliliği için yumuşak veri.", "📐"),
        ("17", "Soru MR", "Soru çözümü anındaki deneyim; mikro geri bildirim.", "❓"),
        ("18", "Deneme MR", "Deneme sonrası psiko-duygusal geri bildirim; churn ve moral erken uyarısı.", "📉"),
        ("19", "Kitap Listesi & Kayıt", "Kurum kütüphanesi; kapak ve PDF ile tek tık erişim.", "📚"),
        ("20", "Kaynak Atama", "Kitap ve kaynağı öğrenciye atama; içerik operasyonunu merkezileştirir.", "🔗"),
        ("21", "Test Oluşturucu", "Kuruma özel test üretimi; ölçüm ve içeriği aynı çatıda birleştirir.", "✏"),
    ]
    for i, (num, tit, bod, ic) in enumerate(cards3):
        r, c = divmod(i, 2)
        x, y, w, h = small_cell(c, r)
        add_card(s3, x, y, w, h, num, tit, bod, ic)

    # --- Slayt 4: Test Maker devam, Reçete, Tarama ---
    s4 = prs.slides.add_slide(blank)
    add_slide_bg(s4)
    add_decor_circles(s4)
    add_accent_bar(s4)
    add_category_title(s4, "05D  İÇERİK, REÇETE, TARAMA", "Üretim Hızı ve Hata Döngüsünü Kapatma")

    add_card(
        s4,
        col0,
        top0,
        big_w,
        big_h,
        "22",
        "Test Maker & Hata Reçetesi",
        "Otomatik Soru Kırpıcı: PDF/görselden soru görüntüsü ayıklama.\nSoru Havuzu: banka ve yeniden kullanım.\n\n"
        "Hata Reçetesi: Reçete Yaz, Reçete Deposu, Hatalı Soru Havuzu ile hedefli tekrar ve kişiselleştirme.",
        "🧩",
    )

    cards4 = [
        ("23", "Otomatik Soru Kırpıcı", "Operasyon süresini kısaltan kırpma hattı; ölçek için kritik.", "✂"),
        ("24", "Soru Havuzu", "Merkezi soru bankası; Test Maker ile entegre içerik varlığı.", "🗂"),
        ("25", "Reçete Yaz / Depo / Havuz", "Koç reçetesi yazımı, şablon arşivi ve hatalı soru havuzu — tek kartta özet.", "💊"),
        ("26", "Tarama Analiz & Rapor", "Tarama sonuçlarının analizi ve raporlanması.", "📑"),
        ("27", "Tarama Oluşturma", "Yeni tarama tanımı ve süreç başlatma.", "➕"),
        ("28", "Tarama Deposu", "Tarama arşivi ve yeniden kullanım.", "🗄"),
    ]
    for i, (num, tit, bod, ic) in enumerate(cards4):
        r, c = divmod(i, 2)
        x, y, w, h = small_cell(c, r)
        add_card(s4, x, y, w, h, num, tit, bod, ic)

    # --- Slayt 5: Yol haritası ---
    s5 = prs.slides.add_slide(blank)
    add_slide_bg(s5)
    add_decor_circles(s5)
    add_accent_bar(s5)
    add_category_title(s5, "05E  YOL HARİTASI", "Şeffaf Ürün Sınırı")
    add_card(
        s5,
        m,
        top0,
        SLIDE_W - 2 * m,
        Inches(2.4),
        "—",
        "Fasikül Oluşturma & Fasikül Deposu",
        "Sidebar’da “Yakında” olarak işaretli; yatırımcıya MVP ile tam vizyon arası mesafenin kontrollü anlatımı. "
        "Pazarlama Asistanı ve Analiz Merkezi BETA etiketiyle erken erişim modunda.",
        "🚀",
    )

    prs.save(path)
    print("Saved:", path)


if __name__ == "__main__":
    import os

    root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    out = os.path.join(root, "docs", "Derecepanel_Koc_Sidebar_Ozellikleri_2026.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    build_presentation(out)
